#!/usr/bin/env python3
"""Acceptance test for the AX1 image-deck assembler and checker."""

from __future__ import annotations

import os
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SKILL_SCRIPTS = ROOT / "skills" / "ax1-presentation" / "scripts"


def run(arguments: list[str], *, expected: int = 0) -> subprocess.CompletedProcess[str]:
    result = subprocess.run(
        [sys.executable, *arguments],
        cwd=ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
        env={**os.environ, "PYTHONIOENCODING": "utf-8"},
    )
    if result.returncode != expected:
        raise AssertionError(
            f"unexpected exit {result.returncode}, expected {expected}\n"
            f"stdout:\n{result.stdout}\nstderr:\n{result.stderr}"
        )
    return result


def main() -> int:
    assembler = str(SKILL_SCRIPTS / "assemble_image_deck.py")
    checker = str(SKILL_SCRIPTS / "check_presentation.py")
    with tempfile.TemporaryDirectory(prefix="ax1-presentation-") as temporary:
        base = Path(temporary)
        project = base / "project"
        project.mkdir()
        images = project / "images"
        images.mkdir()
        Image.new("RGB", (1600, 900), "#10243d").save(images / "slide_01.png")
        Image.new("RGB", (1600, 900), "#e9f1f7").save(images / "slide_02.png")
        (project / "speech.md").write_text(
            "## 슬라이드 1\n\n핵심 메시지를 설명합니다.\n\n"
            "## Slide 2\n\n다음 행동과 예상 질문을 설명합니다.\n",
            encoding="utf-8",
        )

        common = [
            assembler,
            "--project-root",
            str(project),
            "--image-dir",
            "images",
            "--speech",
            "speech.md",
            "--output",
            "candidate.pptx",
        ]
        run(common)
        run(
            [
                checker,
                "--project-root",
                str(project),
                "--pptx",
                "candidate.pptx",
                "--expected-slides",
                "2",
                "--require-notes",
                "--require-full-slide-images",
            ]
        )

        presentation = Presentation(str(project / "candidate.pptx"))
        if "핵심 메시지" not in presentation.slides[0].notes_slide.notes_text_frame.text:
            raise AssertionError("실제 한글 발표자 노트가 PPTX에 보존되지 않음")

        overwrite = run(common, expected=2)
        if "덮어쓰지 않음" not in overwrite.stdout:
            raise AssertionError("기존 출력 덮어쓰기 거부를 확인하지 못함")

        gap_images = project / "gap-images"
        gap_images.mkdir()
        Image.new("RGB", (1600, 900), "white").save(gap_images / "slide_01.png")
        Image.new("RGB", (1600, 900), "white").save(gap_images / "slide_03.png")
        gap = run(
            [
                assembler,
                "--project-root",
                str(project),
                "--image-dir",
                "gap-images",
                "--output",
                "gap.pptx",
            ],
            expected=2,
        )
        if "누락" not in gap.stdout or (project / "gap.pptx").exists():
            raise AssertionError("번호 누락을 출력 없이 거부하지 못함")

        outside_images = base / "outside-images"
        outside_images.mkdir()
        Image.new("RGB", (1600, 900), "white").save(outside_images / "slide_01.png")
        escaped = run(
            [
                assembler,
                "--project-root",
                str(project),
                "--image-dir",
                str(outside_images),
                "--output",
                "escaped.pptx",
            ],
            expected=2,
        )
        if "프로젝트 밖" not in escaped.stdout or (project / "escaped.pptx").exists():
            raise AssertionError("프로젝트 밖 입력을 출력 없이 거부하지 못함")

        wrong_ratio = project / "wrong-ratio"
        wrong_ratio.mkdir()
        Image.new("RGB", (1000, 1000), "white").save(wrong_ratio / "slide_01.png")
        ratio = run(
            [
                assembler,
                "--project-root",
                str(project),
                "--image-dir",
                "wrong-ratio",
                "--output",
                "wrong-ratio.pptx",
            ],
            expected=2,
        )
        if "화면비" not in ratio.stdout or (project / "wrong-ratio.pptx").exists():
            raise AssertionError("잘못된 화면비를 출력 없이 거부하지 못함")

        corrupt_images = project / "corrupt-images"
        corrupt_images.mkdir()
        (corrupt_images / "slide_01.png").write_bytes(b"not-an-image")
        corrupt = run(
            [
                assembler,
                "--project-root",
                str(project),
                "--image-dir",
                "corrupt-images",
                "--output",
                "corrupt.pptx",
            ],
            expected=2,
        )
        if "손상" not in corrupt.stdout or (project / "corrupt.pptx").exists():
            raise AssertionError("손상 이미지를 출력 없이 거부하지 못함")

    print("AX1 presentation acceptance test passed")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
