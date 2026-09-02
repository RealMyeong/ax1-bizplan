#!/usr/bin/env python3
"""Assemble sequential full-slide images into a local PPTX with speaker notes."""

from __future__ import annotations

import argparse
import json
import os
import re
import tempfile
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


IMAGE_PATTERN = re.compile(r"^slide[_-]?(\d+)\.(png|jpe?g)$", re.IGNORECASE)
NOTE_HEADING = re.compile(
    r"^\s{0,3}#{1,6}\s*(?:(?:slide|슬라이드)\s*0*(\d+)|제\s*0*(\d+)\s*(?:페이지|쪽|장))\s*$",
    re.IGNORECASE,
)
MAX_SLIDES = 200
MAX_IMAGE_BYTES = 100 * 1024 * 1024
ASPECTS = {
    "16:9": (13.333333, 7.5),
    "4:3": (10.0, 7.5),
}


class DeckError(ValueError):
    """Raised when the requested deck cannot be assembled safely."""


def project_path(project_root: Path, value: str, *, must_exist: bool = False) -> Path:
    root = project_root.resolve(strict=True)
    raw = Path(value)
    candidate = raw if raw.is_absolute() else root / raw
    resolved = candidate.resolve(strict=must_exist)
    if not resolved.is_relative_to(root):
        raise DeckError(f"프로젝트 밖의 경로는 사용할 수 없음: {value}")
    return resolved


def discover_images(image_dir: Path, expected_ratio: float) -> list[Path]:
    if not image_dir.is_dir():
        raise DeckError(f"이미지 폴더를 찾을 수 없음: {image_dir}")

    numbered: dict[int, Path] = {}
    for path in image_dir.iterdir():
        if not path.is_file():
            continue
        match = IMAGE_PATTERN.fullmatch(path.name)
        if not match:
            continue
        number = int(match.group(1))
        if number < 1:
            raise DeckError(f"슬라이드 번호는 1부터 시작해야 함: {path.name}")
        if number in numbered:
            raise DeckError(
                f"같은 슬라이드 번호가 중복됨: {number} ({numbered[number].name}, {path.name})"
            )
        numbered[number] = path

    if not numbered:
        raise DeckError("slide_01.png 형식의 슬라이드 이미지를 찾지 못함")
    if len(numbered) > MAX_SLIDES:
        raise DeckError(f"최대 슬라이드 수 {MAX_SLIDES}장을 초과함")

    expected_numbers = list(range(1, max(numbered) + 1))
    if sorted(numbered) != expected_numbers:
        missing = sorted(set(expected_numbers) - set(numbered))
        raise DeckError(f"슬라이드 번호가 연속적이지 않음. 누락: {missing}")

    paths = [numbered[number] for number in expected_numbers]
    dimensions: set[tuple[int, int]] = set()
    for path in paths:
        if path.stat().st_size > MAX_IMAGE_BYTES:
            raise DeckError(f"이미지 파일이 100MB를 초과함: {path.name}")
        try:
            with Image.open(path) as image:
                image.verify()
            with Image.open(path) as image:
                width, height = image.size
        except Exception as exc:  # Pillow reports format-specific exceptions.
            raise DeckError(f"손상되었거나 지원하지 않는 이미지: {path.name}: {exc}") from exc
        if width <= 0 or height <= 0:
            raise DeckError(f"유효하지 않은 이미지 크기: {path.name}")
        ratio = width / height
        if abs(ratio - expected_ratio) / expected_ratio > 0.015:
            raise DeckError(
                f"화면비가 맞지 않음: {path.name}={width}x{height}, 기대={expected_ratio:.4f}"
            )
        dimensions.add((width, height))
    if len(dimensions) != 1:
        raise DeckError(f"슬라이드 이미지 크기가 서로 다름: {sorted(dimensions)}")
    return paths


def parse_speaker_notes(path: Path | None) -> dict[int, str]:
    if path is None:
        return {}
    text = path.read_text(encoding="utf-8-sig")
    notes: dict[int, list[str]] = {}
    current: int | None = None
    for line in text.splitlines():
        match = NOTE_HEADING.fullmatch(line)
        if match:
            number = int(match.group(1) or match.group(2))
            if number in notes:
                raise DeckError(f"발표자 노트의 슬라이드 번호가 중복됨: {number}")
            notes[number] = []
            current = number
            continue
        if current is not None:
            notes[current].append(line)
    return {number: "\n".join(lines).strip() for number, lines in notes.items()}


def validate_saved_package(path: Path, slide_count: int) -> None:
    with zipfile.ZipFile(path) as archive:
        bad_member = archive.testzip()
        if bad_member is not None:
            raise DeckError(f"PPTX ZIP CRC 검사 실패: {bad_member}")
    reopened = Presentation(str(path))
    if len(reopened.slides) != slide_count:
        raise DeckError(
            f"PPTX 재개방 후 슬라이드 수 불일치: {len(reopened.slides)} != {slide_count}"
        )


def assemble(
    *,
    project_root: Path,
    image_dir_value: str,
    speech_value: str | None,
    output_value: str,
    aspect: str,
) -> dict[str, object]:
    root = project_root.resolve(strict=True)
    image_dir = project_path(root, image_dir_value, must_exist=True)
    speech = project_path(root, speech_value, must_exist=True) if speech_value else None
    output = project_path(root, output_value, must_exist=False)
    if output.suffix.lower() != ".pptx":
        raise DeckError("출력 파일 확장자는 .pptx여야 함")
    if output.exists():
        raise DeckError(f"기존 출력 파일을 덮어쓰지 않음: {output}")
    if speech is not None and not speech.is_file():
        raise DeckError(f"발표자 노트 파일을 찾을 수 없음: {speech}")

    width_inches, height_inches = ASPECTS[aspect]
    images = discover_images(image_dir, width_inches / height_inches)
    notes = parse_speaker_notes(speech)
    invalid_notes = sorted(number for number in notes if number < 1 or number > len(images))
    if invalid_notes:
        raise DeckError(f"존재하지 않는 슬라이드의 노트가 있음: {invalid_notes}")

    presentation = Presentation()
    presentation.slide_width = Inches(width_inches)
    presentation.slide_height = Inches(height_inches)
    blank_layout = presentation.slide_layouts[6]
    for number, image_path in enumerate(images, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
        note_text = notes.get(number, "")
        if note_text:
            slide.notes_slide.notes_text_frame.text = note_text

    output.parent.mkdir(parents=True, exist_ok=True)
    descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{output.stem}.", suffix=".tmp.pptx", dir=output.parent
    )
    os.close(descriptor)
    temporary = Path(temporary_name)
    try:
        presentation.save(str(temporary))
        validate_saved_package(temporary, len(images))
        os.replace(temporary, output)
    finally:
        if temporary.exists():
            temporary.unlink()

    return {
        "status": "ok",
        "output": output.relative_to(root).as_posix(),
        "slides": len(images),
        "slides_with_notes": sum(bool(notes.get(number)) for number in range(1, len(images) + 1)),
        "aspect": aspect,
        "editable_slide_objects": False,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--project-root", required=True)
    parser.add_argument("--image-dir", required=True)
    parser.add_argument("--speech")
    parser.add_argument("--output", required=True)
    parser.add_argument("--aspect", choices=tuple(ASPECTS), default="16:9")
    args = parser.parse_args()
    try:
        result = assemble(
            project_root=Path(args.project_root),
            image_dir_value=args.image_dir,
            speech_value=args.speech,
            output_value=args.output,
            aspect=args.aspect,
        )
    except (DeckError, FileNotFoundError, OSError, zipfile.BadZipFile) as exc:
        print(json.dumps({"status": "error", "message": str(exc)}, ensure_ascii=False))
        return 2
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
