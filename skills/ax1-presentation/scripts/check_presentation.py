#!/usr/bin/env python3
"""Perform structural checks on a local PPTX without network or office automation."""

from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class CheckError(ValueError):
    """Raised when a presentation fails a required structural check."""


def project_file(project_root: Path, value: str) -> Path:
    root = project_root.resolve(strict=True)
    raw = Path(value)
    candidate = raw if raw.is_absolute() else root / raw
    resolved = candidate.resolve(strict=True)
    if not resolved.is_relative_to(root):
        raise CheckError(f"프로젝트 밖의 파일은 검사할 수 없음: {value}")
    if not resolved.is_file():
        raise CheckError(f"PPTX 파일을 찾을 수 없음: {resolved}")
    return resolved


def inspect(
    *,
    project_root: Path,
    pptx_value: str,
    expected_slides: int | None,
    require_notes: bool,
    require_full_slide_images: bool,
) -> dict[str, object]:
    root = project_root.resolve(strict=True)
    path = project_file(root, pptx_value)
    if path.suffix.lower() != ".pptx":
        raise CheckError("검사 대상 확장자는 .pptx여야 함")
    with zipfile.ZipFile(path) as archive:
        bad_member = archive.testzip()
        if bad_member is not None:
            raise CheckError(f"PPTX ZIP CRC 검사 실패: {bad_member}")

    presentation = Presentation(str(path))
    slide_count = len(presentation.slides)
    if slide_count == 0:
        raise CheckError("슬라이드가 없음")
    if expected_slides is not None and slide_count != expected_slides:
        raise CheckError(f"슬라이드 수 불일치: {slide_count} != {expected_slides}")

    width = presentation.slide_width
    height = presentation.slide_height
    notes_missing: list[int] = []
    full_slide_errors: list[int] = []
    for number, slide in enumerate(presentation.slides, start=1):
        if require_notes:
            note_text = slide.notes_slide.notes_text_frame.text.strip()
            if not note_text:
                notes_missing.append(number)
        if require_full_slide_images:
            pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
            exact = (
                len(pictures) == 1
                and len(slide.shapes) == 1
                and pictures[0].left == 0
                and pictures[0].top == 0
                and pictures[0].width == width
                and pictures[0].height == height
            )
            if not exact:
                full_slide_errors.append(number)

    if notes_missing:
        raise CheckError(f"발표자 노트가 비어 있는 슬라이드: {notes_missing}")
    if full_slide_errors:
        raise CheckError(f"전체 화면 이미지 한 장 구조가 아닌 슬라이드: {full_slide_errors}")

    return {
        "status": "ok",
        "pptx": path.relative_to(root).as_posix(),
        "slides": slide_count,
        "aspect_ratio": round(width / height, 6),
        "notes_required": require_notes,
        "full_slide_images_required": require_full_slide_images,
        "visual_review": "required-separately",
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--project-root", required=True)
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--expected-slides", type=int)
    parser.add_argument("--require-notes", action="store_true")
    parser.add_argument("--require-full-slide-images", action="store_true")
    args = parser.parse_args()
    if args.expected_slides is not None and args.expected_slides < 1:
        parser.error("--expected-slides는 1 이상이어야 함")
    try:
        result = inspect(
            project_root=Path(args.project_root),
            pptx_value=args.pptx,
            expected_slides=args.expected_slides,
            require_notes=args.require_notes,
            require_full_slide_images=args.require_full_slide_images,
        )
    except (CheckError, FileNotFoundError, OSError, zipfile.BadZipFile) as exc:
        print(json.dumps({"status": "error", "message": str(exc)}, ensure_ascii=False))
        return 2
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
